from ascii_colors import ASCIIColors, trace_exception
from sklearn.feature_extraction.text import TfidfVectorizer
import numpy as np
from pathlib import Path
import json
import re
import subprocess
import gc

class NumpyEncoderDecoder(json.JSONEncoder):
    def default(self, obj):
        if isinstance(obj, np.ndarray):
            return {'__numpy_array__': True, 'data': obj.tolist()}
        return super(NumpyEncoderDecoder, self).default(obj)

    @staticmethod
    def as_numpy_array(dct):
        if '__numpy_array__' in dct:
            return np.array(dct['data'])
        return dct


def git_pull(folder_path):
    try:
        # Change the current working directory to the desired folder
        subprocess.run(["git", "checkout", folder_path], check=True, cwd=folder_path)
        # Run 'git pull' in the specified folder
        subprocess.run(["git", "pull"], check=True, cwd=folder_path)
        print("Git pull successful in", folder_path)
    except subprocess.CalledProcessError as e:
        print("Error occurred while executing Git pull:", e)
        # Handle any specific error handling here if required

class AdvancedGarbageCollector:
    @staticmethod
    def hardCollect(obj):
        all_referrers = gc.get_referrers(obj)
        for referrer in all_referrers:
            if not isinstance(referrer, (list, tuple, dict, set)):
                referrer = None
        del obj

    @staticmethod
    def safeHardCollect(variable_name, instance=None):
        if instance is not None:
            if hasattr(instance, variable_name):
                obj = getattr(instance, variable_name)
                AdvancedGarbageCollector.hardCollect(obj)
            else:
                print(f"The variable '{variable_name}' does not exist in the instance.")
        else:
            if variable_name in locals():
                obj = locals()[variable_name]
                AdvancedGarbageCollector.hardCollect(obj)
            elif variable_name in globals():
                obj = globals()[variable_name]
                AdvancedGarbageCollector.hardCollect(obj)
            else:
                print(f"The variable '{variable_name}' does not exist in the local or global namespace.")

    @staticmethod
    def safeHardCollectMultiple(variable_names, instance=None):
        for variable_name in variable_names:
            AdvancedGarbageCollector.safeHardCollect(variable_name, instance)

    @staticmethod
    def collect():
        gc.collect()


class PackageManager:
    @staticmethod
    def install_package(package_name):
        import subprocess
        import sys
        subprocess.check_call([sys.executable, "-m", "pip", "install", package_name])

class GitManager:
    @staticmethod
    def git_pull(folder_path):
        try:
            # Change the current working directory to the desired folder
            subprocess.run(["git", "checkout", folder_path], check=True, cwd=folder_path)
            # Run 'git pull' in the specified folder
            subprocess.run(["git", "pull"], check=True, cwd=folder_path)
            print("Git pull successful in", folder_path)
        except subprocess.CalledProcessError as e:
            print("Error occurred while executing Git pull:", e)
            # Handle any specific error handling here if required

class File64BitsManager:

    @staticmethod
    def raw_b64_img(image) -> str:
        try:
            from PIL import Image, PngImagePlugin
            import io
            import base64
        except:
            PackageManager.install_package("pillow")
            from PIL import Image
            import io
            import base64

        # XXX controlnet only accepts RAW base64 without headers
        with io.BytesIO() as output_bytes:
            metadata = None
            for key, value in image.info.items():
                if isinstance(key, str) and isinstance(value, str):
                    if metadata is None:
                        metadata = PngImagePlugin.PngInfo()
                    metadata.add_text(key, value)
            image.save(output_bytes, format="PNG", pnginfo=metadata)

            bytes_data = output_bytes.getvalue()

        return str(base64.b64encode(bytes_data), "utf-8")


    @staticmethod
    def img2b64(image) -> str:
        return "data:image/png;base64," + File64BitsManager.raw_b64_img(image)    

    @staticmethod
    def b642img(b64img) -> str:
        try:
            from PIL import Image, PngImagePlugin
            import io
            import base64
        except:
            PackageManager.install_package("pillow")
            from PIL import Image
            import io
            import base64        
        image_data = re.sub('^data:image/.+;base64,', '', b64img)
        return Image.open(io.BytesIO(base64.b64decode(image_data)))  

    @staticmethod
    def get_supported_file_extensions_from_base64(b64data):
        # Extract the file extension from the base64 data
        data_match = re.match(r'^data:(.*?);base64,', b64data)
        if data_match:
            mime_type = data_match.group(1)
            extension = mime_type.split('/')[-1]
            return extension
        else:
            raise ValueError("Invalid base64 data format.")
        
    @staticmethod
    def extract_content_from_base64(b64data):
        # Split the base64 data at the comma separator
        header, content = b64data.split(',', 1)

        # Extract only the content part and remove any white spaces and newlines
        content = content.strip()

        return content

    @staticmethod
    def b642file(b64data, filename):
        import base64   
        # Extract the file extension from the base64 data
        
        
        # Save the file with the determined extension
        with open(filename, 'wb') as file:
            file.write(base64.b64decode(File64BitsManager.extract_content_from_base64(b64data)))

        return filename
class TFIDFLoader:
    @staticmethod
    def create_vectorizer_from_dict(tfidf_info):
        vectorizer = TfidfVectorizer(**tfidf_info['params'])
        vectorizer.vocabulary_ = tfidf_info['vocabulary']
        vectorizer.idf_ = [tfidf_info['idf_values'][feature] for feature in vectorizer.get_feature_names()]
        return vectorizer

    @staticmethod
    def create_dict_from_vectorizer(vectorizer):
        tfidf_info = {
            "vocabulary": vectorizer.vocabulary_,
            "idf_values": dict(zip(vectorizer.get_feature_names(), vectorizer.idf_)),
            "params": vectorizer.get_params()
        }
        return tfidf_info
    
class DocumentDecomposer:
    @staticmethod
    def clean_text(text):
        # Remove extra returns and leading/trailing spaces
        text = text.replace('\r', '').strip()
        return text

    @staticmethod
    def split_into_paragraphs(text):
        # Split the text into paragraphs using two or more consecutive newlines
        paragraphs = [p+"\n" for p in re.split(r'\n{2,}', text)]
        return paragraphs

    @staticmethod
    def tokenize_sentences(paragraph):
        # Custom sentence tokenizer using simple regex-based approach
        sentences = [s+"." for s in paragraph.split(".")]
        sentences = [sentence.strip() for sentence in sentences if sentence.strip()]
        return sentences

    @staticmethod
    def decompose_document(text, max_chunk_size, overlap_size, tokenize, detokenize):
        cleaned_text = DocumentDecomposer.clean_text(text)
        paragraphs = DocumentDecomposer.split_into_paragraphs(cleaned_text)

        # List to store the final clean chunks
        clean_chunks = []

        current_chunk = []  # To store the current chunk being built
        l=0
        for paragraph in paragraphs:
            # Tokenize the paragraph into sentences
            sentences = DocumentDecomposer.tokenize_sentences(paragraph)

            for sentence in sentences:
                # If adding the current sentence to the chunk exceeds the max_chunk_size,
                # we add the current chunk to the list of clean chunks and start a new chunk
                tokens = tokenize(sentence)
                nb_tokens = len(tokens)
                if nb_tokens>max_chunk_size:
                    while nb_tokens>max_chunk_size:
                        current_chunk += tokens[:max_chunk_size-l-1]
                        clean_chunks.append(current_chunk)
                        tokens = tokens[max_chunk_size-l-1-overlap_size:]
                        nb_tokens -= max_chunk_size-l-1-overlap_size
                        l=0
                        current_chunk = current_chunk[-overlap_size:]
                else:
                    if l + nb_tokens + 1 > max_chunk_size:

                        clean_chunks.append(current_chunk)
                        if overlap_size==0:
                            current_chunk = []
                        else:
                            current_chunk = current_chunk[-overlap_size:]
                        l=0

                    # Add the current sentence to the chunk
                    current_chunk += tokens
                    l += nb_tokens

        # Add the remaining chunk from the paragraph to the clean_chunks
        if current_chunk:
            clean_chunks.append(current_chunk)
            current_chunk = ""

        return clean_chunks
    
class TextVectorizer:
    def __init__(
                    self, 
                    vectorization_method, # supported "model_embedding" or "ftidf_vectorizer"
                    model=None, #needed in case of using model_embedding
                    database_path=None,
                    save_db=False,
                    visualize_data_at_startup=False,
                    visualize_data_at_add_file=False,
                    visualize_data_at_generate=False,
                    data_visualization_method="PCA",
                    database_dict=None
                    ):
        
        self.vectorization_method = vectorization_method
        self.save_db = save_db
        self.model = model
        self.database_file = database_path

        self.visualize_data_at_startup=visualize_data_at_startup
        self.visualize_data_at_add_file=visualize_data_at_add_file
        self.visualize_data_at_generate=visualize_data_at_generate
        
        self.data_visualization_method = data_visualization_method
        
        if database_dict is not None:
            self.chunks =  database_dict["chunks"]
            self.vectorizer = database_dict["vectorizer"]
            self.infos =   database_dict["infos"]
            self.ready = True
        else:
            self.chunks = {}
            self.ready = False
            self.vectorizer = None
        
            if vectorization_method=="model_embedding":
                try:
                    if not self.model or self.model.embed("hi")==None: # test
                        self.vectorization_method="ftidf_vectorizer"
                        self.infos={
                            "vectorization_method":"ftidf_vectorizer"
                        }
                    else:
                        self.infos={
                            "vectorization_method":"model_embedding"
                        }
                except Exception as ex:
                    ASCIIColors.error("Couldn't embed the text, so trying to use tfidf instead.")
                    trace_exception(ex)
                    self.infos={
                        "vectorization_method":"ftidf_vectorizer"
                    }
        # Load previous state from the JSON file
        if self.save_db:
            if Path(self.database_file).exists():
                ASCIIColors.success(f"Database file found : {self.database_file}")
                self.load_from_json()
                if self.visualize_data_at_startup:
                    self.show_document()
                self.ready = True
            else:
                ASCIIColors.info(f"No database file found : {self.database_file}")

                
    def show_document(self, query_text=None, save_fig_path =None, show_interactive_form=False):
        import textwrap
        import seaborn as sns
        import matplotlib.pyplot as plt
        import mplcursors
        from tkinter import Tk, Text, Scrollbar, Frame, Label, TOP, BOTH, RIGHT, LEFT, Y, N, END

        
        from sklearn.manifold import TSNE
        from sklearn.decomposition import PCA

        if self.data_visualization_method=="PCA":
            use_pca =  True
        else:
            use_pca =  False
        
        if use_pca:
            print("Showing pca representation :")
        else:
            print("Showing t-sne representation :")

        embeddings = {key:chunk["embeddings"] for key, chunk in self.chunks.items()}
        emb = list(embeddings.values())
        ref = list(embeddings.keys())
        if len(emb)>=2:
            # Normalize embeddings
            emb = np.vstack(emb)
            norms = np.linalg.norm(emb, axis=1)
            normalized_embeddings = emb / norms[:, np.newaxis]

            # Embed the query text
            if query_text is not None:
                query_embedding = self.embed_query(query_text)
                query_embedding = query_embedding.detach().squeeze().numpy()
                query_normalized_embedding = query_embedding / np.linalg.norm(query_embedding)

                # Combine the query embeddings with the document embeddings
                combined_embeddings = np.vstack((normalized_embeddings, query_normalized_embedding))
                ref.append("Quey_chunk_0")
            else:
                # Combine the query embeddings with the document embeddings
                combined_embeddings = normalized_embeddings

            if use_pca:
                # Use PCA for dimensionality reduction
                pca = PCA(n_components=2)
                try:
                    embeddings_2d = pca.fit_transform(combined_embeddings)
                except Exception as ex:
                    
                    embeddings_2d = []
            else:
                # Use t-SNE for dimensionality reduction
                # Adjust the perplexity value
                perplexity = min(30, combined_embeddings.shape[0] - 1)
                tsne = TSNE(n_components=2, perplexity=perplexity)
                embeddings_2d = tsne.fit_transform(combined_embeddings)

            # Create a dictionary to map document paths to colors
            document_path_colors = {}
            for i, path in enumerate(ref):
                document_path = "_".join(path.split("_")[:-1])  # Extract the document path (excluding chunk and chunk number)
                if document_path not in document_path_colors:
                    # Assign a new color to the document path if it's not in the dictionary
                    document_path_colors[document_path] = sns.color_palette("hls", len(document_path_colors) + 1)[-1]

            # Generate a list of colors for each data point based on the document path
            point_colors = [document_path_colors["_".join(path.split("_")[:-1])] for path in ref]
            

            # Create a scatter plot using Seaborn
            sns.scatterplot(x=embeddings_2d[:, 0], y=embeddings_2d[:, 1], hue=point_colors)  # Plot document embeddings
            # Add labels to the scatter plot
            for i, (x, y) in enumerate(embeddings_2d[:-1]):
                plt.text(x, y, str(i), fontsize=8)

            plt.xlabel('Dimension 1')
            plt.ylabel('Dimension 2')
            if use_pca:      
                plt.title('Embeddings Scatter Plot based on PCA')
            else:
                plt.title('Embeddings Scatter Plot based on t-SNE')
            # Enable mplcursors to show tooltips on hover
            cursor = mplcursors.cursor(hover=True)

            # Define the hover event handler
            @cursor.connect("add")
            def on_hover(sel):
                index = sel.target.index
                if index > 0:
                    text = self.chunks[index]["chunk_text"]
                    wrapped_text = textwrap.fill(text, width=50)  # Wrap the text into multiple lines
                    sel.annotation.set_text(f"Index: {index}\nText:\n{wrapped_text}")
                else:
                    sel.annotation.set_text("Query")

            # Define the click event handler using matplotlib event handling mechanism
            def on_click(event):
                if event.xdata is not None and event.ydata is not None:
                    x, y = event.xdata, event.ydata
                    distances = ((embeddings_2d[:, 0] - x) ** 2 + (embeddings_2d[:, 1] - y) ** 2)
                    index = distances.argmin()
                    text = self.chunks[index]["chunk_text"] if index < len(self.chunks) else query_text

                    # Open a new Tkinter window with the content of the text
                    root = Tk()
                    root.title(f"Text for Index {index}")
                    frame = Frame(root)
                    frame.pack(fill=BOTH, expand=True)

                    label = Label(frame, text="Text:")
                    label.pack(side=TOP, padx=5, pady=5)

                    text_box = Text(frame)
                    text_box.pack(side=TOP, padx=5, pady=5, fill=BOTH, expand=True)
                    text_box.insert(END, text)

                    scrollbar = Scrollbar(frame)
                    scrollbar.pack(side=RIGHT, fill=Y)
                    scrollbar.config(command=text_box.yview)
                    text_box.config(yscrollcommand=scrollbar.set)

                    text_box.config(state="disabled")

                    root.mainloop()

            # Connect the click event handler to the figure
            plt.gcf().canvas.mpl_connect("button_press_event", on_click)
            if save_fig_path:
                try:
                    plt.savefig(save_fig_path)
                except Exception as ex:
                    trace_exception(ex)
            if show_interactive_form:
                plt.show()
    
    def file_exists(self, document_name:str)->bool:
        # Loop through the list of dictionaries
        for dictionary in self.chunks:
            if 'document_name' in dictionary and dictionary['document_name'] == document_name:
                # If the document_name is found in the current dictionary, set the flag to True and break the loop
                document_name_found = True
                return True
        return False
    
    def remove_document(self, document_name:str):
        for dictionary in self.chunks:
            if 'document_name' in dictionary and dictionary['document_name'] == document_name:
                # If the document_name is found in the current dictionary, set the flag to True and break the loop
                self.chunks.remove(dictionary)
                return True
        return False



    def add_document(self, document_name:Path, text:str, chunk_size: int, overlap_size:int, force_vectorize=False,add_as_a_bloc=False):
        if self.file_exists(document_name) and not force_vectorize:
            print(f"Document {document_name} already exists. Skipping vectorization.")
            return
        if add_as_a_bloc:
            chunks_text = [self.model.tokenize(text)]
            for i, chunk in enumerate(chunks_text):
                chunk_id = f"{document_name}_chunk_{i + 1}"
                chunk_dict = {
                    "document_name": document_name,
                    "chunk_index": i+1,
                    "chunk_text":self.model.detokenize(chunk),
                    "chunk_tokens": chunk,
                    "embeddings":[]
                }
                self.chunks[chunk_id] = chunk_dict
        else:
            chunks_text = DocumentDecomposer.decompose_document(text, chunk_size, overlap_size, self.model.tokenize, self.model.detokenize)
            for i, chunk in enumerate(chunks_text):
                chunk_id = f"{document_name}_chunk_{i + 1}"
                chunk_dict = {
                    "document_name": document_name,
                    "chunk_index": i+1,
                    "chunk_text":self.model.detokenize(chunk),
                    "chunk_tokens": chunk,
                    "embeddings":[]
                }
                self.chunks[chunk_id] = chunk_dict
        
    def index(self):
        if self.vectorization_method=="ftidf_vectorizer":
            self.vectorizer = TfidfVectorizer()
            #if self.debug:
            #    ASCIIColors.yellow(','.join([len(chunk) for chunk in chunks]))
            data=[]
            for k,chunk in self.chunks.items():
                try:
                    data.append(chunk["chunk_text"]) 
                except Exception as ex:
                    print("oups")
            self.vectorizer.fit(data)

        # Generate embeddings for each chunk
        for chunk_id, chunk in self.chunks.items():
            # Store chunk ID, embeddings, and original text
            try:
                if self.vectorization_method=="ftidf_vectorizer":
                    chunk["embeddings"] = self.vectorizer.transform([chunk["chunk_text"]]).toarray()
                else:
                    chunk["embeddings"] = self.model.embed(chunk["chunk_text"])
            except Exception as ex:
                print("oups")

        if self.save_db:
            self.save_to_json()
            
        self.ready = True
        if self.visualize_data_at_add_file:
            self.show_document()


    def embed_query(self, query_text):
        # Generate query embeddings
        if self.vectorization_method=="ftidf_vectorizer":
            query_embedding = self.vectorizer.transform([query_text]).toarray()
        else:
            query_embedding = self.model.embed(query_text)
            if query_embedding is None:
                ASCIIColors.warning("The model doesn't implement embeddings extraction")
                self.vectorization_method="ftidf_vectorizer"
                query_embedding = self.vectorizer.transform([query_text]).toarray()

        return query_embedding

    def recover_text(self, query_embedding, top_k=1):
        from sklearn.metrics.pairwise import cosine_similarity
        similarities = {}
        for chunk_id, chunk in self.chunks.items():
            similarity = cosine_similarity(query_embedding, chunk["embeddings"])
            similarities[chunk_id] = similarity

        # Sort the similarities and retrieve the top-k most similar embeddings
        sorted_similarities = sorted(similarities.items(), key=lambda x: x[1], reverse=True)[:top_k]

        # Retrieve the original text associated with the most similar embeddings
        texts = [self.chunks[chunk_id]["chunk_text"] for chunk_id, _ in sorted_similarities]

        if self.visualize_data_at_generate:
            self.show_document()

        return texts, sorted_similarities

    def toJson(self):
        state = {
            "chunks": self.chunks,
            "infos": self.infos,
            "vectorizer": TFIDFLoader.create_vectorizer_from_dict(self.vectorizer) if self.vectorization_method=="ftidf_vectorizer" else None
        }
        return state
    
    def setVectorizer(self, vectorizer_dict:dict):
        self.vectorizer=TFIDFLoader.create_vectorizer_from_dict(vectorizer_dict)

    def save_to_json(self):
        state = {
            "chunks": self.chunks,
            "infos": self.infos,
            "vectorizer": TFIDFLoader.create_vectorizer_from_dict(self.vectorizer) if self.vectorization_method=="ftidf_vectorizer" else None
        }
        with open(self.database_file, "w") as f:
            json.dump(state, f, cls=NumpyEncoderDecoder, indent=4)

    def load_from_json(self):

        ASCIIColors.info("Loading vectorized documents")
        with open(self.database_file, "r") as f:
            database = json.load(f, object_hook=NumpyEncoderDecoder.as_numpy_array)
            self.chunks = database["chunks"]
            self.infos= database["infos"]
            self.ready = True
        if self.vectorization_method=="ftidf_vectorizer":
            from sklearn.feature_extraction.text import TfidfVectorizer
            data = [c["chunk_text"] for k,c in self.chunks.items()]
            if len(data)>0:
                self.vectorizer = TfidfVectorizer()
                self.vectorizer.fit(data)
                self.embeddings={}
                for k,chunk in self.chunks.items():
                    chunk["embeddings"][k]= self.vectorizer.transform([chunk["embeddings"]]).toarray()
                    
                    
    def clear_database(self):
        self.ready = False
        self.vectorizer=None
        self.chunks = {}
        self.infos={}
        if self.save_db:
            self.save_to_json()
            
      
class GenericDataLoader:
    @staticmethod        
    def read_file(file_path:Path):
        if file_path.suffix ==".pdf":
            return GenericDataLoader.read_pdf_file(file_path)
        elif file_path.suffix == ".docx":
            return GenericDataLoader.read_docx_file(file_path)
        elif file_path.suffix == ".json":
            return GenericDataLoader.read_json_file(file_path)
        elif file_path.suffix == ".html":
            return GenericDataLoader.read_html_file(file_path)
        elif file_path.suffix == ".pptx":
            return GenericDataLoader.read_pptx_file(file_path)
        if file_path.suffix in [".txt", ".rtf", ".md", ".log", ".cpp", ".java", ".js", ".py", ".rb", ".sh", ".sql", ".css", ".html", ".php", ".json", ".xml", ".yaml", ".yml", ".h", ".hh", ".hpp", ".inc", ".snippet", ".snippets", ".asm", ".s", ".se", ".sym", ".ini", ".inf", ".map", ".bat"]:
            return GenericDataLoader.read_text_file(file_path)
        else:
            raise ValueError("Unknown file type")
    def get_supported_file_types():
        return ["pdf", "txt", "docx", "json", "html", "pptx",".txt", ".md", ".log", ".cpp", ".java", ".js", ".py", ".rb", ".sh", ".sql", ".css", ".html", ".php", ".json", ".xml", ".yaml", ".yml", ".h", ".hh", ".hpp", ".inc", ".snippet", ".snippets", ".asm", ".s", ".se", ".sym", ".ini", ".inf", ".map", ".bat", ".rtf"]    
    @staticmethod        
    def read_pdf_file(file_path):
        try:
            import PyPDF2
            from PIL import Image, UnidentifiedImageError
            import pytesseract
            import pdfminer
            from pdfminer.high_level import extract_text
        except ImportError:
            PackageManager.install_package("PyPDF2")
            PackageManager.install_package("pytesseract")
            PackageManager.install_package("pillow")
            PackageManager.install_package("pdfminer")
            PackageManager.install_package("pdfminer.six")
            
            import PyPDF2
            from PIL import Image, UnidentifiedImageError
            import pytesseract
            from pdfminer.high_level import extract_text
            
        # Extract text from the PDF
        text = extract_text(file_path)

        # Convert to Markdown (You may need to implement custom logic based on your specific use case)
        markdown_text = text.replace('\n', '  \n')  # Adding double spaces at the end of each line for Markdown line breaks
        
        return markdown_text

    @staticmethod
    def read_docx_file(file_path):
        try:
            from docx import Document
        except ImportError:
            PackageManager.install_package("python-docx")
            from docx import Document
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text

    @staticmethod
    def read_json_file(file_path):
        import json
        with open(file_path, 'r') as file:
            data = json.load(file)
        return data
    
    @staticmethod
    def read_csv_file(file_path):
        try:
            import csv
        except ImportError:
            PackageManager.install_package("csv")
            import csv
        with open(file_path, 'r') as file:
            csv_reader = csv.reader(file)
            lines = [row for row in csv_reader]
        return lines    

    @staticmethod
    def read_html_file(file_path):
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            PackageManager.install_package("beautifulsoup4")
            from bs4 import BeautifulSoup
        with open(file_path, 'r') as file:
            soup = BeautifulSoup(file, 'html.parser')
            text = soup.get_text()
        return text
    
    @staticmethod
    def read_pptx_file(file_path):
        try:
            from pptx import Presentation
        except ImportError:
            PackageManager.install_package("python-pptx")
            from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
        return text
    
    @staticmethod
    def read_text_file(file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        return content




class PromptReshaper:
    def __init__(self, template:str):
        self.template = template
    def replace(self, placeholders:dict)->str:
        template = self.template
        # Calculate the number of tokens for each placeholder
        for placeholder, text in placeholders.items():
            template = template.replace(placeholder, text)
        return template
    def build(self, placeholders:dict, tokenize, detokenize, max_nb_tokens:int, place_holders_to_sacrifice:list=[])->str:
        # Tokenize the template without placeholders
        template_text = self.template
        for placeholder in placeholders:
            template_text = template_text.replace("{{" + placeholder + "}}", "")
        template_tokens = tokenize(template_text)
        
        # Calculate the number of tokens in the template without placeholders
        template_tokens_count = len(template_tokens)
        
        # Calculate the number of tokens for each placeholder
        placeholder_tokens_count = {}
        all_count = template_tokens_count
        for placeholder, text in placeholders.items():
            text_tokens = tokenize(text)
            placeholder_tokens_count[placeholder] = len(text_tokens)
            all_count += placeholder_tokens_count[placeholder]

        def fill_template(template, data):
            for key, value in data.items():
                placeholder = "{{" + key + "}}"
                template = template.replace(placeholder, value)
            return template
        
        if max_nb_tokens-all_count>0 or len(place_holders_to_sacrifice)==0:
            return fill_template(self.template, placeholders)
        else:
            to_remove = -int((max_nb_tokens - all_count)/len(place_holders_to_sacrifice))
            for placeholder, text in placeholders.items():
                if placeholder in place_holders_to_sacrifice:
                    text_tokens = tokenize(text)[to_remove:]
                    placeholders[placeholder]=detokenize(text_tokens)
            return fill_template(self.template, placeholders)



class LOLLMSLocalizer:
    def __init__(self, dictionary):
        self.dictionary = dictionary

    def localize(self, input_string):
        def replace(match):
            key = match.group(1)
            return self.dictionary.get(key, match.group(0))
        
        import re
        pattern = r'@<([^>]+)>@'
        localized_string = re.sub(pattern, replace, input_string)
        return localized_string


class File_Path_Generator:
    @staticmethod
    def generate_unique_file_path(folder_path, file_base_name, file_extension):
        folder_path = Path(folder_path)
        index = 0
        while True:
            # Construct the full file path with the current index
            file_name = f"{file_base_name}_{index}.{file_extension}"
            full_file_path = folder_path / file_name
            
            # Check if the file already exists in the folder
            if not full_file_path.exists():
                return full_file_path
            
            # If the file exists, increment the index and try again
            index += 1